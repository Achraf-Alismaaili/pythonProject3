import pandas as pd


import plotly.express as px
import os
from pptx import Presentation
import streamlit as st
from io import BytesIO
import io
import ftplib
import tempfile
from PIL import Image
def add_logo(logo_path, width, height):
    """Read and return a resized logo"""
    logo = Image.open(logo_path)
    modified_logo = logo.resize((width, height))
    return modified_logo
my_logo = add_logo(logo_path="OCP_Group.svg.png", width=200, height=200)

st.set_page_config(page_title="Dashboard", page_icon=my_logo, layout="wide")
st.sidebar.image(my_logo)
@st.cache
def get_data_from_excel(file):
    df = pd.read_excel(
        io=file,
        engine="openpyxl",
    )
    return df
def plot_to_pic(f):
    buf = io.BytesIO()
    f.write_image(buf, format="jpg")
    fp = tempfile.NamedTemporaryFile()
    with open(f"{fp.name}.jpg", 'wb') as ff:
        ff.write(buf.getvalue())
    return buf, f"{fp.name}.jpg"

uploaded_file = st.file_uploader("Choose a file",type=['xlsx'])
if uploaded_file is not None:
    df =get_data_from_excel(uploaded_file)

# par mois


    df2=pd.DataFrame()
    df2['mois']=df['mois'].unique()
    df2['JHF_prévu']=[0]*len(df2)
    df2['JHF_réalisé']=[0]*len(df2)
    c=0
    for e in df2['mois']:
        df2.iat[c,1]=df[df['mois']==e]['DUREE'].mean()*df[df['mois']==e].count()['PRESENCE']
        df2.iat[c,2]=df[df['mois']==e][df['PRESENCE']=='Présent']['DUREE'].mean()*df[df['mois']==e][df['PRESENCE']=='Présent'].count()['PRESENCE']
        c=c+1
    df2['Taux']=df2['JHF_réalisé']/df2['JHF_prévu']*100
    df2=df2.round(2)

    inter = pd.DataFrame()
    l = df2['mois'].unique()
    months = [x for pair in zip(l, l) for x in pair]
    inter['mois'] = months
    inter['type_JHF'] = ['prévu', 'réalisé'] * len(l)
    val = []
    for i in range(len(inter['type_JHF'])):
        if inter.iat[i, 1] == 'prévu':
            val.append(float(df2[df2['mois'] == inter.iat[i, 0]]['JHF_prévu']))
        else:
            val.append(float(df2[df2['mois'] == inter.iat[i, 0]]['JHF_réalisé']))

    inter['valeur'] = val



    fig = px.bar(inter, x="mois", color="type_JHF",
                 y='valeur',
                 barmode='group',
                 title='<b>JHF prévu/JHF réalisé</b>',
                 text='valeur'

                 )

    fig.update_layout(
        plot_bgcolor="rgba(0,0,0,0)",
        xaxis=(dict(showgrid=False))
    )


    #### par entité


    df3=pd.DataFrame()
    df3['Entité2']=df['Entité2'].unique()
    df3['JHF_prévu']=[0]*len(df3)
    df3['JHF_réalisé']=[0]*len(df3)
    c=0
    for e in df3['Entité2']:
        df3.iat[c,1]=df[df['Entité2']==e]['DUREE'].mean()*df[df['Entité2']==e].count()['PRESENCE']
        df3.iat[c,2]=df[df['Entité2']==e]['DUREE'].mean()*df[df['Entité2']==e][df['PRESENCE']=='Présent'].count()['PRESENCE']
        c=c+1
    df3['Taux']=df3['JHF_réalisé']/df3['JHF_prévu']*100
    df3=df3.round(2)
    df3 = df3.sort_values(by=['JHF_prévu'], ascending=False)
    inter3 = pd.DataFrame()
    l3 = df3['Entité2'].unique()
    ent = [x for pair in zip(l3, l3) for x in pair]
    inter3['Entité2'] = ent
    inter3['type_JHF'] = ['prévu', 'réalisé'] * len(l3)
    val = []
    for i in range(len(inter3['type_JHF'])):
        if inter3.iat[i, 1] == 'prévu':
            val.append(float(df3[df3['Entité2'] == inter3.iat[i, 0]]['JHF_prévu']))
        else:
            val.append(float(df3[df3['Entité2'] == inter3.iat[i, 0]]['JHF_réalisé']))

    inter3['valeur'] = val


    fig3 = px.bar(inter3, x="Entité2", color="type_JHF",
                  y='valeur',
                  barmode='group',
                  title="<b>JHF réalise par entité</b>")

    fig3.update_layout(
        plot_bgcolor="rgba(0,0,0,0)",
        xaxis=(dict(showgrid=False))
    )


    # par categorie





    df4 = pd.DataFrame()

    df4['CATEGORIE'] = df['CATEGORIE'].unique()




    list4 = df.groupby(['CATEGORIE']).mean()['DUREE'].tolist()
    a4 = df.count()['PRESENCE']
    df4['JHF_prévu'] = [element * a4 for element in list4]




    list4 = df.groupby(['CATEGORIE']).mean()['DUREE'].tolist()
    b4 = df[df['PRESENCE'] == 'Présent'].count()['PRESENCE']
    df4['JHF_réalisé'] = [element * b4 for element in list4]


    df4['Taux'] = df4['JHF_réalisé'] / df4['JHF_prévu'] * 100

    df4 = df4.round(2)

    fig4 = px.pie(df4, names='CATEGORIE', values='JHF_réalisé', title='<b>Répartition des JHF par catégorie</b>')
    fig4.update_traces(textposition='inside', textinfo='percent+label')
    fig4.update_layout(
        plot_bgcolor="rgba(0,0,0,0)",
        xaxis=(dict(showgrid=False))
    )



    # par type de formation



    df5=pd.DataFrame()
    df5['MACRO PROCESS']=df['MACRO PROCESS'].unique()
    df5['JHF_prévu']=[0]*len(df5)
    df5['JHF_réalisé']=[0]*len(df5)
    c=0
    for e in df5['MACRO PROCESS']:
        df5.iat[c,1]=df[df['MACRO PROCESS']==e]['DUREE'].mean()*df[df['MACRO PROCESS']==e].count()['PRESENCE']
        df5.iat[c,2]=df[df['MACRO PROCESS']==e]['DUREE'].mean()*df[df['MACRO PROCESS']==e][df['PRESENCE']=='Présent'].count()['PRESENCE']
        c=c+1
    df5['Taux']=df5['JHF_réalisé']/df5['JHF_prévu']*100
    df5=df5.round(2)
    #df5.sort_values(by=['JHF_réalisé'], ascending=False)



    fig5 = px.bar(df5.sort_values(by=['JHF_réalisé']), y='MACRO PROCESS', x='JHF_réalisé', orientation='h',
                  text='JHF_réalisé', title='<b>Répartition JHF par type de formation</b>')
    fig5.update_layout(
        plot_bgcolor="rgba(0,0,0,0)",
        xaxis=(dict(showgrid=False))
    )






    ##effectif inscrit/prsent





    df6 = pd.DataFrame()



    df6['mois'] = df['mois'].unique()


    df6['effectif_inscrit'] = [0] * len(df6)

    df6['effectif_présent'] = [0] * len(df6)

    c = 0
    for e in df6['mois']:
        df6.iat[c, 1] = df[df['mois'] == e]['NOM_PRENOM2'].count()
        df6.iat[c, 2] = df[df['mois'] == e][df['PRESENCE'] == 'Présent']['NOM_PRENOM2'].count()
        c = c + 1
    inter6 = pd.DataFrame()
    l6 = df6['mois'].unique()
    months = [x for pair in zip(l6, l6) for x in pair]
    inter6['mois'] = months
    inter6['type_effectif'] = ['inscrit', 'présent'] * len(l6)
    val = []
    for i in range(len(inter6['type_effectif'])):
        if inter6.iat[i, 1] == 'inscrit':
            val.append(float(df6[df6['mois'] == inter6.iat[i, 0]]['effectif_inscrit']))
        else:
            val.append(float(df6[df6['mois'] == inter6.iat[i, 0]]['effectif_présent']))

    inter6['valeur'] = val


    fig6 = px.bar(inter6, x="mois", color="type_effectif",
                  y='valeur',
                  barmode='group', text='valeur',
                  title='<b>effectif prévu/effectif réalisé</b>',

                  )

    fig6.update_layout(
        plot_bgcolor="rgba(0,0,0,0)",
        xaxis=(dict(showgrid=False))
    )



















    ################
    #st.dataframe(df_selection)
    st.title(":bar_chart: tableau de bord_version1")
    st.markdown("---")


    column1,column2= st.columns(2)
    column1.plotly_chart(fig, use_container_width=True)
    column2.plotly_chart(fig6, use_container_width=True)





    column11,column3,column22 = st.columns(3)
    column11.plotly_chart(fig5, use_container_width=True)
    column3.plotly_chart(fig4, use_container_width=True)
    column22.plotly_chart(fig3, use_container_width=True)





    b1,name1=plot_to_pic(fig)

    b2, name2 = plot_to_pic(fig6)

    b3, name3 = plot_to_pic(fig5)

    b4, name4 = plot_to_pic(fig4)

    b5, name5 = plot_to_pic(fig3)






    binary_output = BytesIO()
    prs = Presentation()
    blank_slide_layout=prs.slide_layouts[1]

    images=[name1,name2,name3,name4,name5]
    bufs=[b1,b2,b3,b4,b5]
    for i in range(len(images)):
        slide = prs.slides.add_slide(blank_slide_layout)
        picture = slide.shapes.add_picture(images[i], 0, 0)


    #picture = slide.shapes.add_picture("images/fig3.png", 0, 0, width=4, height=4)




    prs.save(binary_output)

    btn=st.download_button(label = 'Download powerpoint',
                       data = binary_output.getvalue(),
                       #on_click=delI(),
                       file_name = 'my_power.pptx')

    output = BytesIO()
    writer = pd.ExcelWriter(output, engine='xlsxwriter')
    df2.to_excel(writer, index=False, sheet_name='Sheet1')
    workbook = writer.book
    worksheet = writer.sheets['Sheet1']
    format1 = workbook.add_format({'num_format': '0.00'})
    worksheet.set_column('A:A', None, format1)
    worksheet.insert_image('F4',name1)
    b1.close()

    df6.to_excel(writer, index=False, sheet_name='Sheet2')
    worksheet2 = writer.sheets['Sheet2']
    worksheet2.set_column('A:A', None, format1)
    worksheet2.insert_image('F4',name2)
    b2.close()


    df5.to_excel(writer, index=False, sheet_name='Sheet3')
    worksheet3 = writer.sheets['Sheet3']
    worksheet3.set_column('A:A', None, format1)
    worksheet3.insert_image('F4',name3)
    b3.close()

    df4.to_excel(writer, index=False, sheet_name='Sheet4')
    worksheet4 = writer.sheets['Sheet4']
    worksheet4.set_column('A:A', None, format1)
    worksheet4.insert_image('F4',name4)
    b4.close()

    df3.to_excel(writer, index=False, sheet_name='Sheet5')
    worksheet5 = writer.sheets['Sheet5']
    worksheet5.set_column('A:A', None, format1)
    worksheet5.insert_image('F4',name5)
    b5.close()


    writer.save()
    processed_data = output.getvalue()
    st.download_button(label='📥 Download back up',
                                    data=processed_data ,
                                    file_name= 'backup_test.xlsx')

    hide_st_style = """
                <style>
                #MainMenu {visibility: hidden;}
                footer {visibility: hidden;}
                header {visibility: hidden;}
                </style>
                """
    st.markdown(hide_st_style, unsafe_allow_html=True)

